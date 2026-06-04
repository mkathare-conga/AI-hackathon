"""
Generate final hackathon pitch deck — 4 live agents.
Includes: team members, business problem, expected ROI.
Output: docs/pitch_deck.pptx
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─── Palette ──────────────────────────────────────────────────────────────────
TEAL       = RGBColor(0x16, 0x60, 0x5A)
WARM_BG    = RGBColor(0xF5, 0xF1, 0xE8)
DARK       = RGBColor(0x1A, 0x22, 0x24)
MUTED      = RGBColor(0x5A, 0x6A, 0x72)
ALERT      = RGBColor(0x8B, 0x2F, 0x1D)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_TEAL = RGBColor(0xE2, 0xF0, 0xEF)
GREEN      = RGBColor(0x27, 0xAE, 0x60)

# ─── Setup ────────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

TOTAL = 10


# ─── Helpers ──────────────────────────────────────────────────────────────────
def blank_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])


def bg(slide, color=WARM_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def text(slide, content, left, top, width, height,
         size=Pt(18), bold=False, color=DARK, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = content
    r.font.size = size
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return txBox


def bullets(slide, items, left, top, width, height, size=Pt(15), color=DARK, dot_color=TEAL):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(6)
        p.clear()
        dot_run = p.add_run()
        dot_run.text = "● "
        dot_run.font.size = Pt(8)
        dot_run.font.color.rgb = dot_color
        dot_run.font.bold = True
        dot_run.font.name = "Calibri"
        txt_run = p.add_run()
        txt_run.text = item
        txt_run.font.size = size
        txt_run.font.color.rgb = color
        txt_run.font.name = "Calibri"
    return txBox


def pill(slide, label, left, top, fill=LIGHT_TEAL, color=TEAL):
    b = box(slide, left, top, 3.0, 0.32, fill_color=fill)
    tf = b.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label.upper()
    r.font.size = Pt(9)
    r.font.bold = True
    r.font.color.rgb = color
    r.font.name = "Calibri"


def teal_bar(slide):
    box(slide, 0, 0, 13.33, 0.06, fill_color=TEAL)


def snum(slide, n):
    text(slide, f"{n}/{TOTAL}", left=12.4, top=7.1, width=0.8, height=0.3,
         size=Pt(9), color=MUTED, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title + Team
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
bg(s)
teal_bar(s)
box(s, 0, 0.06, 5.2, 7.44, fill_color=TEAL)

text(s, "Commercial Execution\nIntelligence Platform",
     left=0.45, top=1.4, width=4.4, height=2.4,
     size=Pt(32), bold=True, color=WHITE)

text(s, "4 Live AI Agents\nPre & Post-Signature Revenue Protection",
     left=0.45, top=3.8, width=4.2, height=1.0,
     size=Pt(16), color=RGBColor(0xC8, 0xE6, 0xE4))

text(s, "Team Members",
     left=0.45, top=5.2, width=4.2, height=0.4,
     size=Pt(10), bold=True, color=RGBColor(0xA0, 0xD0, 0xCC))

text(s, "Vighneshwar Bhat\nRaghavendra Rao\nManoj Kathare",
     left=0.45, top=5.55, width=4.2, height=1.2,
     size=Pt(14), color=WHITE)

# Right side
text(s, "Hackathon Final Demo · June 2026",
     left=5.7, top=2.0, width=6.0, height=0.5,
     size=Pt(12), bold=True, color=MUTED)

text(s, "AI-powered platform that optimizes pre-signature\npricing and detects post-signature revenue leakage,\ncontract drift, and amendment impact.",
     left=5.7, top=2.7, width=6.5, height=1.5,
     size=Pt(18), color=DARK)

# Key numbers
text(s, "$54,720", left=5.7, top=4.8, width=2.5, height=0.6,
     size=Pt(26), bold=True, color=ALERT)
text(s, "missed revenue found", left=5.7, top=5.35, width=2.5, height=0.3,
     size=Pt(10), color=MUTED)

text(s, "$570K", left=8.3, top=4.8, width=2.0, height=0.6,
     size=Pt(26), bold=True, color=ALERT)
text(s, "contract drift impact", left=8.3, top=5.35, width=2.3, height=0.3,
     size=Pt(10), color=MUTED)

text(s, "+$194K", left=10.6, top=4.8, width=2.0, height=0.6,
     size=Pt(26), bold=True, color=GREEN)
text(s, "amendment impact", left=10.6, top=5.35, width=2.3, height=0.3,
     size=Pt(10), color=MUTED)

snum(s, 1)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Business Problem
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
bg(s)
teal_bar(s)

pill(s, "The Business Problem", left=0.6, top=0.4)
text(s, "Revenue leaks after signature when\ncontract truth is not operationalized.",
     left=0.6, top=0.88, width=9.0, height=1.5,
     size=Pt(30), bold=True, color=DARK)

text(s, "Enterprise companies lose millions annually because pricing rights, uplift clauses, and renewal\nobligations agreed in contracts are never enforced in downstream billing and operational systems.",
     left=0.6, top=2.5, width=9.5, height=0.8, size=Pt(15), color=MUTED)

bullets(s, [
    "Contract says 5% annual uplift — billing stays flat year after year",
    "Sales quotes one price — contract gets signed at a different rate",
    "Amendment changes terms — no one updates billing or provisioning",
    "Renewal notice window passes — the revenue uplift right expires silently",
    "No single system connects contract truth to operational reality",
    "Pre-signature pricing lacks data — sellers leave money on the table",
], left=0.6, top=3.4, width=8.2, height=3.2, size=Pt(15))

# Right side impact
b = box(s, 9.4, 1.2, 3.5, 5.5, fill_color=LIGHT_TEAL, line_color=TEAL, line_width=Pt(1.5))
text(s, "Industry Impact",
     left=9.55, top=1.35, width=3.2, height=0.35, size=Pt(10), bold=True, color=TEAL)
text(s, "9.7%\nof billings contain\npricing errors",
     left=9.55, top=1.9, width=3.2, height=1.2, size=Pt(20), bold=True, color=ALERT, align=PP_ALIGN.CENTER)
text(s, "42%\nof companies have\nno audit process",
     left=9.55, top=3.3, width=3.2, height=1.2, size=Pt(20), bold=True, color=ALERT, align=PP_ALIGN.CENTER)
text(s, "$3M+\naverage annual\nleakage per $100M\nin contract value",
     left=9.55, top=4.7, width=3.2, height=1.5, size=Pt(18), bold=True, color=ALERT, align=PP_ALIGN.CENTER)

snum(s, 2)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Platform Architecture
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
bg(s)
teal_bar(s)

pill(s, "Our Solution", left=0.6, top=0.4)
text(s, "One shared evidence layer.\nFour agents on top.",
     left=0.6, top=0.88, width=9.0, height=1.4,
     size=Pt(30), bold=True, color=DARK)

# Evidence layer
box(s, 0.8, 2.55, 11.7, 0.85, fill_color=TEAL)
text(s, "Shared Evidence Layer  ·  Document Ingestion  ·  AI Clause Extraction  ·  Governing-Term Resolution",
     left=0.9, top=2.62, width=11.5, height=0.7,
     size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Four agent cards
agents = [
    ("Revenue Leakage\nInvestigator", "LIVE", "$54,720 found"),
    ("Quote-to-Contract\nDrift Detector", "LIVE", "$570K drift impact"),
    ("Amendment Impact\nDetector", "LIVE", "+$194K net delta"),
    ("Pre-Sign Pricing\nAdvisor", "LIVE", "Pre-signature optimization"),
]
card_w = 2.7
gap = 0.2
start_x = 0.8
for i, (name, status, stats) in enumerate(agents):
    x = start_x + i * (card_w + gap)
    box(s, x, 3.6, card_w, 2.2, fill_color=LIGHT_TEAL, line_color=TEAL, line_width=Pt(1.5))
    text(s, name, left=x+0.1, top=3.7, width=card_w-0.2, height=0.8,
         size=Pt(13), bold=True, color=DARK, align=PP_ALIGN.CENTER)
    text(s, f"● {status}", left=x+0.1, top=4.45, width=card_w-0.2, height=0.3,
         size=Pt(10), bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    text(s, stats, left=x+0.1, top=4.85, width=card_w-0.2, height=0.6,
         size=Pt(11), color=MUTED, align=PP_ALIGN.CENTER)

# Data sources below
text(s, "↑", left=6.4, top=6.0, width=0.5, height=0.4,
     size=Pt(18), bold=True, color=MUTED, align=PP_ALIGN.CENTER)
text(s, "Contracts (PDF/DOCX)  ·  Amendments  ·  Invoices/ERP  ·  Quotes  ·  Renewal Events  ·  Historical Deals",
     left=0.8, top=6.4, width=11.7, height=0.4,
     size=Pt(12), color=MUTED, align=PP_ALIGN.CENTER)

snum(s, 3)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Agent 1: Revenue Leakage Investigator
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
bg(s)
teal_bar(s)

pill(s, "Agent 1 — Revenue Leakage", left=0.6, top=0.4)
text(s, "Revenue Leakage\nInvestigator",
     left=0.6, top=0.88, width=8.0, height=1.4,
     size=Pt(32), bold=True, color=TEAL)

text(s, "Finds missed revenue the business was already entitled to charge.",
     left=0.6, top=2.4, width=8.5, height=0.5, size=Pt(16), color=MUTED)

text(s, "WHAT IT DETECTS", left=0.6, top=3.1, width=3.0, height=0.3,
     size=Pt(9), bold=True, color=TEAL)
bullets(s, [
    "Missed uplift — invoice vs. contract obligation comparison",
    "At-risk accounts — notice deadline approaching without action",
    "Amendment precedence — resolves which doc legally controls the rate",
    "Dollar quantification — exact per-account missed revenue",
], left=0.6, top=3.4, width=6.5, height=2.4, size=Pt(14))

# Demo callout
b = box(s, 8.8, 2.9, 4.0, 3.6, fill_color=LIGHT_TEAL, line_color=TEAL, line_width=Pt(1.2))
text(s, "LIVE RESULTS", left=8.95, top=3.0, width=3.7, height=0.3,
     size=Pt(9), bold=True, color=TEAL)
text(s, "3 leakage cases\n$54,720 total missed\n1 at-risk prediction\n$6,000 at stake",
     left=8.95, top=3.4, width=3.7, height=1.5, size=Pt(14), color=DARK)
text(s, "Key insight:",
     left=8.95, top=5.0, width=3.7, height=0.3, size=Pt(10), bold=True, color=ALERT)
text(s, "Resolves conflicting contract\nterms to find the single\ncontrolling source of truth.",
     left=8.95, top=5.3, width=3.7, height=1.0, size=Pt(12), color=ALERT)

snum(s, 4)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Agent 2: Drift Detector
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
bg(s)
teal_bar(s)

pill(s, "Agent 2 — Drift Detector", left=0.6, top=0.4)
text(s, "Quote-to-Contract\nDrift Detector",
     left=0.6, top=0.88, width=8.0, height=1.4,
     size=Pt(32), bold=True, color=TEAL)

text(s, "Compares what was quoted vs. what was signed — catches negotiation drift.",
     left=0.6, top=2.4, width=9.0, height=0.5, size=Pt(16), color=MUTED)

text(s, "HOW IT WORKS", left=0.6, top=3.1, width=3.0, height=0.3,
     size=Pt(9), bold=True, color=TEAL)
bullets(s, [
    "AI extracts contract facts from uploaded documents",
    "Line-by-line comparison: price, quantity, term, SLA per quote line",
    "Scope drift detection: items quoted but missing from contract",
    "Impact scoring with severity classification (high/medium/low)",
], left=0.6, top=3.4, width=6.5, height=2.4, size=Pt(14))

# Results callout
b = box(s, 8.8, 2.9, 4.0, 3.6, fill_color=LIGHT_TEAL, line_color=TEAL, line_width=Pt(1.2))
text(s, "LIVE RESULTS", left=8.95, top=3.0, width=3.7, height=0.3,
     size=Pt(9), bold=True, color=TEAL)
text(s, "3 quotes analyzed\n11 drift findings\n6 high severity\n$570,360 total impact",
     left=8.95, top=3.4, width=3.7, height=1.5, size=Pt(14), color=DARK)
text(s, "Accounts:",
     left=8.95, top=5.0, width=3.7, height=0.3, size=Pt(10), bold=True, color=MUTED)
text(s, "TechVault: $155K\nMeridian: $249K\nAtlas: $165K",
     left=8.95, top=5.3, width=3.7, height=1.0, size=Pt(12), color=DARK)

snum(s, 5)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Agent 3: Amendment Impact
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
bg(s)
teal_bar(s)

pill(s, "Agent 3 — Amendment Impact", left=0.6, top=0.4)
text(s, "Amendment Impact\nDetector",
     left=0.6, top=0.88, width=8.0, height=1.4,
     size=Pt(32), bold=True, color=TEAL)

text(s, "Identifies downstream operational changes triggered by contract amendments.",
     left=0.6, top=2.4, width=9.0, height=0.5, size=Pt(16), color=MUTED)

text(s, "WHAT IT PRODUCES", left=0.6, top=3.1, width=3.0, height=0.3,
     size=Pt(9), bold=True, color=TEAL)
bullets(s, [
    "Impact analysis: what changed, before vs. after, severity",
    "Revenue delta: quantifies annual revenue gain/loss per change",
    "Action items: billing updates, workflow changes, provisioning needed",
    "Team routing: auto-assigns to RevOps, Engineering, Legal, CS",
], left=0.6, top=3.4, width=6.5, height=2.4, size=Pt(14))

# Results callout
b = box(s, 8.8, 2.9, 4.0, 3.8, fill_color=LIGHT_TEAL, line_color=TEAL, line_width=Pt(1.2))
text(s, "LIVE RESULTS", left=8.95, top=3.0, width=3.7, height=0.3,
     size=Pt(9), bold=True, color=TEAL)
text(s, "3 amendments analyzed\n11 impacts detected\n14 open action items\n+$194,400 net delta",
     left=8.95, top=3.4, width=3.7, height=1.5, size=Pt(14), color=DARK)
text(s, "GlobalTech: +$86K ▲\nCascade:    –$72K ▼\nVertex:     +$180K ▲",
     left=8.95, top=5.1, width=3.7, height=1.0, size=Pt(12), color=DARK)

snum(s, 6)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Agent 4: Pricing Advisor
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
bg(s)
teal_bar(s)

pill(s, "Agent 4 — Pre-Signature", left=0.6, top=0.4)
text(s, "Pre-Sign Pricing\nAdvisor",
     left=0.6, top=0.88, width=8.0, height=1.4,
     size=Pt(32), bold=True, color=TEAL)

text(s, "Before signature: data-driven pricing recommendations that maximize revenue while preserving deal confidence.",
     left=0.6, top=2.4, width=9.0, height=0.5, size=Pt(16), color=MUTED)

text(s, "WHAT IT DELIVERS", left=0.6, top=3.1, width=3.0, height=0.3,
     size=Pt(9), bold=True, color=TEAL)
bullets(s, [
    "Recommended improved unit price based on same-company historical deals",
    "Close-confidence tradeoff — shows risk of pushing price higher",
    "Pricing floor calculation — minimum acceptable price per account",
    "Incremental ACV and total contract value impact quantification",
    "Approval path and signing considerations for deal desk",
], left=0.6, top=3.4, width=6.5, height=2.8, size=Pt(14))

# Results callout
b = box(s, 8.8, 2.9, 4.0, 3.8, fill_color=LIGHT_TEAL, line_color=TEAL, line_width=Pt(1.2))
text(s, "LIVE CAPABILITIES", left=8.95, top=3.0, width=3.7, height=0.3,
     size=Pt(9), bold=True, color=TEAL)
text(s, "Historical deal analysis\nSame-company win rates\nConfidence scoring\nPrice optimization",
     left=8.95, top=3.4, width=3.7, height=1.5, size=Pt(14), color=DARK)
text(s, "Value proposition:",
     left=8.95, top=5.1, width=3.7, height=0.3, size=Pt(10), bold=True, color=GREEN)
text(s, "Helps sellers capture\n3-8% more revenue per\ndeal without losing\nclose rate.",
     left=8.95, top=5.4, width=3.7, height=1.0, size=Pt(12), color=GREEN)

snum(s, 7)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Expected Business Value & ROI
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
bg(s)
teal_bar(s)

pill(s, "Business Value & ROI", left=0.6, top=0.4)
text(s, "Measurable revenue impact\nfrom day one.",
     left=0.6, top=0.88, width=9.0, height=1.4,
     size=Pt(30), bold=True, color=DARK)

# ROI boxes — row 1
roi_items = [
    ("$54,720", "Missed Revenue\nRecovered", "Immediate uplift from\nenforcing existing contract\nrights already earned.", ALERT),
    ("$570,360", "Contract Drift\nPrevented", "Quote-to-contract gaps\ncaught before they become\npermanent revenue loss.", ALERT),
    ("+$194,400", "Amendment Impact\nQuantified", "Net positive revenue from\namendments routed to the\nright teams for action.", GREEN),
    ("3-8%", "Pricing Uplift\nPer Deal", "Data-driven pricing raises\nACV without reducing\nclose confidence.", GREEN),
]

card_w = 2.9
gap = 0.18
start_x = 0.6
for i, (val, label, desc, val_color) in enumerate(roi_items):
    x = start_x + i * (card_w + gap)
    box(s, x, 2.5, card_w, 3.3, fill_color=LIGHT_TEAL, line_color=TEAL, line_width=Pt(1.2))
    text(s, val, left=x+0.1, top=2.65, width=card_w-0.2, height=0.6,
         size=Pt(22), bold=True, color=val_color, align=PP_ALIGN.CENTER)
    text(s, label, left=x+0.1, top=3.3, width=card_w-0.2, height=0.7,
         size=Pt(13), bold=True, color=DARK, align=PP_ALIGN.CENTER)
    text(s, desc, left=x+0.15, top=4.1, width=card_w-0.3, height=1.3,
         size=Pt(11), color=MUTED, align=PP_ALIGN.CENTER)

# Bottom summary
text(s, "For a company with $100M in contract value:",
     left=0.6, top=6.1, width=12.0, height=0.4, size=Pt(14), bold=True, color=DARK)
text(s, "Conservative estimate: $800K–$1.2M in recoverable or protectable revenue annually  ·  ROI payback: < 3 months",
     left=0.6, top=6.5, width=12.0, height=0.5, size=Pt(14), color=TEAL)

snum(s, 8)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Design Principles + Tech Stack
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
bg(s)
teal_bar(s)

pill(s, "How We Built It", left=0.6, top=0.4)
text(s, "AI reads. Logic calculates.\nNeither alone is enough.",
     left=0.6, top=0.88, width=9.0, height=1.4,
     size=Pt(28), bold=True, color=DARK)

# Two columns
box(s, 0.6, 2.4, 5.8, 2.8, fill_color=LIGHT_TEAL, line_color=TEAL, line_width=Pt(1))
text(s, "What AI does", left=0.8, top=2.55, width=5.4, height=0.35,
     size=Pt(12), bold=True, color=TEAL)
bullets(s, [
    "Extracts clauses from unstructured contract language",
    "Resolves which document controls a given term",
    "Generates investigation briefs and pricing insights",
    "Handles varied formats, messy language, version drift",
], left=0.8, top=2.95, width=5.4, height=2.0, size=Pt(13), color=DARK, dot_color=TEAL)

box(s, 6.9, 2.4, 5.8, 2.8, fill_color=RGBColor(0xF0, 0xF0, 0xF0),
    line_color=MUTED, line_width=Pt(1))
text(s, "What deterministic logic does", left=7.1, top=2.55, width=5.4, height=0.35,
     size=Pt(12), bold=True, color=MUTED)
bullets(s, [
    "Calculates exact dollar impact — no hallucination",
    "Evaluates notice windows and deadline math",
    "Line-by-line comparison with severity scoring",
    "Reproducible, auditable results every run",
], left=7.1, top=2.95, width=5.4, height=2.0, size=Pt(13), color=DARK, dot_color=MUTED)

# Tech stack summary
text(s, "TECH STACK", left=0.6, top=5.5, width=3.0, height=0.3,
     size=Pt(9), bold=True, color=TEAL)
stack_items = [
    "React + Vite frontend",
    "Python FastAPI backend",
    "GPT-4o via GitHub Models",
    "PostgreSQL 16",
    "MinIO object store",
    "Docker Compose — single command deploy",
]
bullets(s, stack_items, left=0.6, top=5.8, width=12.0, height=1.5, size=Pt(12), color=DARK, dot_color=TEAL)

snum(s, 9)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Close / Thank You
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
bg(s, TEAL)

text(s, "Commercial Execution\nIntelligence Platform",
     left=1.0, top=1.0, width=11.3, height=2.0,
     size=Pt(38), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

text(s, "Turns contract truth into operational action.",
     left=1.0, top=3.0, width=11.3, height=0.7,
     size=Pt(22), color=LIGHT_TEAL, align=PP_ALIGN.CENTER)

# Summary boxes
summaries = [
    ("$54,720", "revenue\nrecovered"),
    ("$570K", "drift\nprevented"),
    ("+$194K", "amendment\nimpact"),
    ("3-8%", "pricing\nuplift"),
]
for i, (val, label) in enumerate(summaries):
    x = 1.2 + i * 3.0
    box(s, x, 3.9, 2.6, 1.5, fill_color=RGBColor(0x1E, 0x7A, 0x72), line_color=LIGHT_TEAL, line_width=Pt(1))
    text(s, val, left=x+0.05, top=4.0, width=2.5, height=0.6,
         size=Pt(22), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    text(s, label, left=x+0.05, top=4.6, width=2.5, height=0.6,
         size=Pt(11), color=RGBColor(0xC8, 0xE6, 0xE4), align=PP_ALIGN.CENTER)

text(s, "4 live agents  ·  Pre & post-signature  ·  AI + deterministic hybrid  ·  Production-ready",
     left=1.0, top=5.7, width=11.3, height=0.5,
     size=Pt(13), color=RGBColor(0xA0, 0xD0, 0xCC), align=PP_ALIGN.CENTER)

text(s, "Team: Vighneshwar Bhat  ·  Raghavendra Rao  ·  Manoj Kathare",
     left=1.0, top=6.3, width=11.3, height=0.4,
     size=Pt(12), color=RGBColor(0x99, 0xC4, 0xC1), align=PP_ALIGN.CENTER)

text(s, "Thank you  ·  Questions?",
     left=1.0, top=6.9, width=11.3, height=0.4,
     size=Pt(16), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

snum(s, 10)


# ─── Save ─────────────────────────────────────────────────────────────────────
out = Path(__file__).resolve().parent.parent / "docs" / "pitch_deck.pptx"
out.parent.mkdir(parents=True, exist_ok=True)
# If file is locked, write to a temp name and replace
try:
    prs.save(out)
except PermissionError:
    tmp = out.with_name("pitch_deck_new.pptx")
    prs.save(tmp)
    print(f"⚠ Original locked — saved as: {tmp}")
    print(f"  Close the open file and rename manually, or re-run.")
    raise SystemExit(0)
print(f"✓ Saved: {out}")
print(f"  {TOTAL} slides — final hackathon pitch deck")
