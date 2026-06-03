"""
Generate Commercial Execution Intelligence Platform deck.
Clean, minimal — no heavy text.
Output: docs/commercial_execution_platform.pptx
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ─── Palette ──────────────────────────────────────────────────────────────────
TEAL       = RGBColor(0x16, 0x60, 0x5A)
WARM_BG    = RGBColor(0xF5, 0xF1, 0xE8)
DARK       = RGBColor(0x1A, 0x22, 0x24)
MUTED      = RGBColor(0x5A, 0x6A, 0x72)
ALERT      = RGBColor(0x8B, 0x2F, 0x1D)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_TEAL = RGBColor(0xE2, 0xF0, 0xEF)

# ─── Slide size: widescreen 16:9 ──────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


# ─── Helpers ──────────────────────────────────────────────────────────────────
def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)


def bg(slide, color=WARM_BG):
    """Fill slide background."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             size=Pt(18), bold=False, color=DARK,
             align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txBox


def add_bullets(slide, items, left, top, width, height,
                size=Pt(16), color=DARK, dot_color=TEAL, line_gap=Pt(8)):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = line_gap
        run = p.add_run()
        run.text = f"  {item}"
        run.font.size = size
        run.font.color.rgb = color
        run.font.name = "Calibri"
        # teal bullet dot (prepend)
        dot = p.runs[0]  # same run — we'll use a prefix char instead
        # re-do with separate runs for dot + text
        p.clear()
        dot_run = p.add_run()
        dot_run.text = "● "
        dot_run.font.size = Pt(8)
        dot_run.font.color.rgb = dot_color
        dot_run.font.bold = True
        dot_run.font.name = "Calibri"
        txt_run = p.add_run()
        txt_run.text = item
        txt_run.font.size = size
        txt_run.font.color.rgb = color
        txt_run.font.name = "Calibri"
    return txBox


def pill_label(slide, text, left, top):
    """Small teal pill tag."""
    w, h = 3.2, 0.32
    b = box(slide, left, top, w, h, fill_color=LIGHT_TEAL)
    tf = b.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text.upper()
    r.font.size = Pt(9)
    r.font.bold = True
    r.font.color.rgb = TEAL
    r.font.name = "Calibri"
    return b


def teal_bar(slide, height=0.06):
    """Thin teal accent bar at top."""
    box(slide, 0, 0, 13.33, height, fill_color=TEAL)


def slide_number(slide, n, total):
    add_text(slide, f"{n} / {total}",
             left=12.4, top=7.1, width=0.8, height=0.3,
             size=Pt(9), color=MUTED, align=PP_ALIGN.RIGHT)


# ─── Slide content ────────────────────────────────────────────────────────────
TOTAL = 8


# ── Slide 1: Title ────────────────────────────────────────────────────────────
s = blank_slide(prs)
bg(s)
teal_bar(s)

# Left teal block
box(s, 0, 0.06, 5.2, 7.44, fill_color=TEAL)

# Title text on teal
add_text(s, "Commercial Execution\nIntelligence Platform",
         left=0.45, top=1.6, width=4.4, height=2.4,
         size=Pt(32), bold=True, color=WHITE)

add_text(s, "Revenue Leakage Investigator\nis the first live agent",
         left=0.45, top=4.2, width=4.2, height=1.2,
         size=Pt(17), bold=False, color=RGBColor(0xC8, 0xE6, 0xE4))

# Right side
add_text(s, "AI-powered platform for post-signature\ncommercial execution",
         left=5.7, top=2.4, width=7.0, height=1.0,
         size=Pt(19), bold=False, color=MUTED)

pill_label(s, "Hackathon Demo · May 2026", left=5.7, top=3.6)

add_text(s, "Conga",
         left=5.7, top=6.9, width=2.0, height=0.4,
         size=Pt(10), color=MUTED)

slide_number(s, 1, TOTAL)


# ── Slide 2: The Problem ───────────────────────────────────────────────────────
s = blank_slide(prs)
bg(s)
teal_bar(s)

pill_label(s, "The Problem", left=0.6, top=0.4)
add_text(s, "The contract is signed.\nRevenue execution drifts.",
         left=0.6, top=0.88, width=8.5, height=1.5,
         size=Pt(30), bold=True, color=DARK)

add_text(s, "Pricing rights, uplift clauses, and renewal obligations\nare agreed in contracts — but rarely operationalized downstream.",
         left=0.6, top=2.55, width=8.5, height=1.0,
         size=Pt(16), color=MUTED)

bullets = [
    "Contract says 5% annual uplift — billing stays flat",
    "Amendment changes terms — ERP keeps using old rates",
    "Renewal notice window passes — revenue uplift right expires",
    "No single system connects contract truth to billing reality",
]
add_bullets(s, bullets, left=0.6, top=3.6, width=8.2, height=2.6, size=Pt(15))

# Right example callout
b = box(s, 9.4, 1.6, 3.5, 4.5, fill_color=LIGHT_TEAL, line_color=TEAL, line_width=Pt(1.5))
add_text(s, "Real example",
         left=9.55, top=1.75, width=3.1, height=0.35,
         size=Pt(10), bold=True, color=TEAL)
add_text(s, "Contract\n5% annual uplift\n30-day notice required",
         left=9.55, top=2.2, width=3.1, height=0.9,
         size=Pt(13), color=DARK)
add_text(s, "↓", left=10.7, top=3.2, width=0.6, height=0.4, size=Pt(18), bold=True, color=ALERT, align=PP_ALIGN.CENTER)
add_text(s, "Billing\nFlat rate. No change.\nEvery month. For years.",
         left=9.55, top=3.65, width=3.1, height=0.9,
         size=Pt(13), color=ALERT)
add_text(s, "Gap = leaked revenue",
         left=9.55, top=4.75, width=3.1, height=0.4,
         size=Pt(12), bold=True, color=ALERT)

slide_number(s, 2, TOTAL)


# ── Slide 3: Platform Story ────────────────────────────────────────────────────
s = blank_slide(prs)
bg(s)
teal_bar(s)

pill_label(s, "The Platform", left=0.6, top=0.4)
add_text(s, "One shared evidence layer.\nMany agents on top.",
         left=0.6, top=0.88, width=9.0, height=1.4,
         size=Pt(30), bold=True, color=DARK)

# Center evidence layer box
box(s, 1.5, 2.55, 10.3, 0.9, fill_color=TEAL)
add_text(s, "Shared Evidence Layer  ·  Document Ingestion  ·  Clause Extraction  ·  Governing-Term Resolution  ·  AI Briefs",
         left=1.55, top=2.62, width=10.2, height=0.75,
         size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Agent cards above (on top visually = below in slide order but positioned above)
agents = [
    ("Revenue Leakage\nInvestigator", "Live ✓", TEAL),
    ("Quote-to-Contract\nDrift Detector", "Planned", MUTED),
    ("Amendment\nImpact Detector", "Planned", MUTED),
    ("Billing vs Contract\nMismatch Finder", "Planned", MUTED),
]
card_w = 2.4
gap = 0.2
start_x = 1.5
for i, (name, status, color) in enumerate(agents):
    x = start_x + i * (card_w + gap)
    fill = LIGHT_TEAL if color == TEAL else RGBColor(0xF0, 0xF0, 0xF0)
    border = TEAL if color == TEAL else RGBColor(0xCC, 0xCC, 0xCC)
    box(s, x, 3.65, card_w, 1.4, fill_color=fill, line_color=border, line_width=Pt(1.2))
    add_text(s, name, left=x+0.1, top=3.75, width=card_w-0.2, height=0.8,
             size=Pt(12), bold=(color == TEAL), color=DARK, align=PP_ALIGN.CENTER)
    status_color = TEAL if status == "Live ✓" else MUTED
    add_text(s, status, left=x+0.1, top=4.55, width=card_w-0.2, height=0.4,
             size=Pt(10), bold=True, color=status_color, align=PP_ALIGN.CENTER)

add_text(s, "Contracts, amendments, invoices, and renewal signals",
         left=1.5, top=5.3, width=10.3, height=0.4,
         size=Pt(12), color=MUTED, align=PP_ALIGN.CENTER)
add_text(s, "↑", left=6.4, top=5.7, width=0.5, height=0.4,
         size=Pt(20), bold=True, color=MUTED, align=PP_ALIGN.CENTER)
add_text(s, "Documents & Billing Data  (PDF, DOCX, ERP extract)",
         left=1.5, top=6.1, width=10.3, height=0.4,
         size=Pt(11), color=MUTED, align=PP_ALIGN.CENTER)

slide_number(s, 3, TOTAL)


# ── Slide 4: Revenue Leakage Investigator ─────────────────────────────────────
s = blank_slide(prs)
bg(s)
teal_bar(s)

pill_label(s, "First Live Agent", left=0.6, top=0.4)
add_text(s, "Revenue Leakage\nInvestigator",
         left=0.6, top=0.88, width=8.0, height=1.5,
         size=Pt(34), bold=True, color=TEAL)

add_text(s, "Finds and quantifies missed revenue that the business was already entitled to charge.",
         left=0.6, top=2.55, width=8.5, height=0.7,
         size=Pt(16), color=MUTED)

bullets_l = [
    "Detects missed uplift — invoice vs. contract comparison",
    "Predicts upcoming notice deadline failures",
    "Resolves which document legally controls the rate",
    "Quantifies exact dollar impact per account",
]
bullets_r = [
    "Multi-document clause extraction (PDF + DOCX)",
    "AI governs-term resolution across conflicting files",
    "Deterministic revenue math — no hallucinations on numbers",
    "Source-of-record verdict with evidence chain",
]
add_bullets(s, bullets_l, left=0.6, top=3.35, width=5.8, height=2.8, size=Pt(14))
add_bullets(s, bullets_r, left=6.9, top=3.35, width=5.8, height=2.8, size=Pt(14))

add_text(s, "WHAT IT DOES", left=0.6, top=3.15, width=3.0, height=0.25,
         size=Pt(9), bold=True, color=TEAL)
add_text(s, "HOW IT WORKS", left=6.9, top=3.15, width=3.0, height=0.25,
         size=Pt(9), bold=True, color=TEAL)

slide_number(s, 4, TOTAL)


# ── Slide 5: Governing Term Resolution ────────────────────────────────────────
s = blank_slide(prs)
bg(s)
teal_bar(s)

pill_label(s, "How It Decides", left=0.6, top=0.4)
add_text(s, "Conflicting documents are resolved,\nnot just listed.",
         left=0.6, top=0.88, width=9.0, height=1.4,
         size=Pt(30), bold=True, color=DARK)

# Precedence chain boxes
chain = [
    ("Amendment", "Explicit override language\nbeats everything", TEAL, WHITE),
    ("Renewal Notice", "Confirms or resets\ncommitted terms", RGBColor(0x2A, 0x8A, 0x82), WHITE),
    ("Order Form", "Deal-specific schedule\noverrides MSA defaults", MUTED, WHITE),
    ("MSA", "Baseline terms —\nlowest precedence", RGBColor(0xCC, 0xCC, 0xCC), DARK),
]
box_w = 2.6
for i, (label, desc, fill, txt) in enumerate(chain):
    x = 0.5 + i * (box_w + 0.25)
    box(s, x, 2.6, box_w, 1.0, fill_color=fill)
    add_text(s, label, left=x+0.1, top=2.65, width=box_w-0.2, height=0.45,
             size=Pt(14), bold=True, color=txt, align=PP_ALIGN.CENTER)
    add_text(s, desc, left=x+0.1, top=3.1, width=box_w-0.2, height=0.55,
             size=Pt(10), color=txt, align=PP_ALIGN.CENTER)
    if i < 3:
        add_text(s, ">", left=x+box_w+0.02, top=2.88, width=0.22, height=0.35,
                 size=Pt(16), bold=True, color=MUTED, align=PP_ALIGN.CENTER)

add_text(s, "Higher precedence →", left=0.5, top=3.75, width=4.0, height=0.35,
         size=Pt(10), color=MUTED)

add_bullets(s, [
    "AI reads every uploaded document and extracts all uplift/override clauses",
    "Ranks by document type, version, and recency",
    "Flags the single controlling term with source evidence",
    "Unrelated documents are attached but do not affect revenue signals",
], left=0.6, top=4.3, width=8.2, height=2.5, size=Pt(14))

# Demo example callout
b = box(s, 9.3, 2.45, 3.7, 4.5, fill_color=LIGHT_TEAL, line_color=TEAL, line_width=Pt(1.2))
add_text(s, "Demo: Summit Distribution",
         left=9.45, top=2.55, width=3.4, height=0.35,
         size=Pt(10), bold=True, color=TEAL)
add_text(s, "MSA v1 → 4% uplift\nAmendment v2 → 6%\nAmendment v3 → 8%\nAmendment v4 → 10% ✓ wins",
         left=9.45, top=3.0, width=3.3, height=1.5,
         size=Pt(12), color=DARK)
add_text(s, "Leakage: $30,720",
         left=9.45, top=4.6, width=3.3, height=0.4,
         size=Pt(13), bold=True, color=ALERT)
add_text(s, "Evidence: Amendment v4\ncontrols. Source verified.",
         left=9.45, top=5.1, width=3.3, height=0.7,
         size=Pt(11), color=MUTED)

slide_number(s, 5, TOTAL)


# ── Slide 6: Why AI + Deterministic ───────────────────────────────────────────
s = blank_slide(prs)
bg(s)
teal_bar(s)

pill_label(s, "Design Principle", left=0.6, top=0.4)
add_text(s, "AI reads. Logic calculates.\nNeither alone is enough.",
         left=0.6, top=0.88, width=9.0, height=1.4,
         size=Pt(30), bold=True, color=DARK)

# Two columns
box(s, 0.6, 2.55, 5.8, 3.8, fill_color=LIGHT_TEAL, line_color=TEAL, line_width=Pt(1))
add_text(s, "What AI does",
         left=0.8, top=2.7, width=5.4, height=0.4,
         size=Pt(13), bold=True, color=TEAL)
add_bullets(s, [
    "Extracts clause meaning from unstructured contract language",
    "Resolves which document currently controls a term",
    "Generates investigation briefs and case explanations",
    "Handles messy language, varied formats, version drift",
], left=0.8, top=3.2, width=5.4, height=2.8, size=Pt(14),
   color=DARK, dot_color=TEAL)

box(s, 6.9, 2.55, 5.8, 3.8, fill_color=RGBColor(0xF0, 0xF0, 0xF0),
    line_color=MUTED, line_width=Pt(1))
add_text(s, "What deterministic logic does",
         left=7.1, top=2.7, width=5.4, height=0.4,
         size=Pt(13), bold=True, color=MUTED)
add_bullets(s, [
    "Calculates exact dollar impact — no approximation",
    "Evaluates notice windows and deadline math",
    "Scores and ranks cases for the investigation queue",
    "Ensures reproducible, auditable results every time",
], left=7.1, top=3.2, width=5.4, height=2.8, size=Pt(14),
   color=DARK, dot_color=MUTED)

add_text(s, "The result: explainable findings backed by source evidence — safe to show a judge or a customer.",
         left=0.6, top=6.5, width=12.1, height=0.6,
         size=Pt(13), color=MUTED, align=PP_ALIGN.CENTER)

slide_number(s, 6, TOTAL)


# ── Slide 7: Roadmap ───────────────────────────────────────────────────────────
s = blank_slide(prs)
bg(s)
teal_bar(s)

pill_label(s, "Roadmap", left=0.6, top=0.4)
add_text(s, "Same platform. More agents.",
         left=0.6, top=0.88, width=9.0, height=1.0,
         size=Pt(32), bold=True, color=DARK)

add_text(s, "No rebuild. The next agent reuses ingestion, extraction, and evidence resolution on day one.",
         left=0.6, top=2.0, width=9.5, height=0.6,
         size=Pt(16), color=MUTED)

roadmap = [
    ("Revenue Leakage\nInvestigator", "LIVE NOW", "Missed uplift and\nrenewal detection", TEAL, "✓"),
    ("Quote-to-Contract\nDrift Detector", "NEXT", "Quoted vs. signed\nterm comparison", MUTED, "○"),
    ("Amendment Impact\nDetector", "PLANNED", "Downstream billing\nchange alerts", MUTED, "○"),
    ("Billing vs Contract\nMismatch Finder", "PLANNED", "Live billing cross-\nreference engine", MUTED, "○"),
]
card_w = 2.8
for i, (name, phase, desc, color, icon) in enumerate(roadmap):
    x = 0.6 + i * (card_w + 0.42)
    fill = TEAL if color == TEAL else RGBColor(0xF2, 0xF2, 0xF2)
    txt_color = WHITE if color == TEAL else DARK
    box(s, x, 2.85, card_w, 3.5, fill_color=fill)
    add_text(s, icon, left=x+0.15, top=2.95, width=0.5, height=0.5,
             size=Pt(18), bold=True, color=WHITE if color == TEAL else MUTED)
    add_text(s, phase, left=x+0.15, top=3.45, width=card_w-0.3, height=0.35,
             size=Pt(9), bold=True,
             color=LIGHT_TEAL if color == TEAL else MUTED)
    add_text(s, name, left=x+0.15, top=3.85, width=card_w-0.3, height=0.8,
             size=Pt(14), bold=True, color=txt_color)
    add_text(s, desc, left=x+0.15, top=4.75, width=card_w-0.3, height=0.8,
             size=Pt(12), color=WHITE if color == TEAL else MUTED)

slide_number(s, 7, TOTAL)


# ── Slide 8: Close ─────────────────────────────────────────────────────────────
s = blank_slide(prs)
bg(s, TEAL)

add_text(s, "Commercial Execution\nIntelligence Platform",
         left=1.0, top=1.5, width=11.3, height=2.2,
         size=Pt(38), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s, "Turns contract truth into operational action.",
         left=1.0, top=3.8, width=11.3, height=0.8,
         size=Pt(22), bold=False, color=LIGHT_TEAL, align=PP_ALIGN.CENTER)

add_text(s, "Revenue Leakage Investigator is the first live agent —\nshowing how AI reconstructs commercial source of truth, detects missed execution,\nand quantifies the financial impact before more revenue slips away.",
         left=2.0, top=4.75, width=9.3, height=1.5,
         size=Pt(14), color=RGBColor(0xC8, 0xE6, 0xE4), align=PP_ALIGN.CENTER)

add_text(s, "Conga  ·  Hackathon 2026",
         left=1.0, top=6.8, width=11.3, height=0.4,
         size=Pt(10), color=RGBColor(0x99, 0xC4, 0xC1), align=PP_ALIGN.CENTER)

slide_number(s, 8, TOTAL)


# ─── Save ─────────────────────────────────────────────────────────────────────
out = Path(__file__).resolve().parent.parent / "docs" / "commercial_execution_platform.pptx"
prs.save(out)
print(f"Saved → {out}")
