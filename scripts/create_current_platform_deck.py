from __future__ import annotations

import json
import os
from pathlib import Path
from urllib.error import HTTPError, URLError
from urllib.request import urlopen

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


API_BASE = os.environ.get("DECK_API_BASE", "http://localhost:8001")
OUTPUT_PATH = Path(__file__).resolve().parents[1] / "docs" / "commercial_execution_platform_updated.pptx"
TOTAL_SLIDES = 9

TEAL = RGBColor(0x16, 0x60, 0x5A)
WARM_BG = RGBColor(0xF5, 0xF1, 0xE8)
DARK = RGBColor(0x1F, 0x1B, 0x16)
MUTED = RGBColor(0x6A, 0x62, 0x56)
ALERT = RGBColor(0x8B, 0x2F, 0x1D)
WARNING = RGBColor(0xB2, 0x67, 0x1F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_TEAL = RGBColor(0xE2, 0xF0, 0xEF)
LIGHT_CARD = RGBColor(0xFF, 0xFC, 0xF5)
GREEN = RGBColor(0x27, 0xAE, 0x60)
LIGHT_BORDER = RGBColor(0xD9, 0xD2, 0xC7)


def fetch_json(path: str):
    url = f"{API_BASE}{path}"
    try:
        with urlopen(url) as response:
            return json.load(response)
    except HTTPError as exc:
        raise RuntimeError(f"Request failed for {path}: {exc.code}") from exc
    except URLError as exc:
        raise RuntimeError(f"Could not reach API at {url}: {exc.reason}") from exc


def ensure_list(payload):
    if isinstance(payload, list):
        return payload
    if isinstance(payload, dict) and "value" in payload and isinstance(payload["value"], list):
        return payload["value"]
    raise RuntimeError("Expected a list payload from the API.")


def fmt_currency(value: float, digits: int = 0, signed: bool = False) -> str:
    prefix = "+" if signed and value > 0 else "-" if signed and value < 0 else ""
    return f"{prefix}${abs(value):,.{digits}f}"


def compact_currency(value: float, signed: bool = False) -> str:
    prefix = "+" if signed and value > 0 else "-" if signed and value < 0 else ""
    amount = abs(value)
    if amount >= 1_000_000:
        rendered = f"${amount / 1_000_000:.2f}M"
    elif amount >= 1_000:
        rendered = f"${amount / 1_000:.1f}K"
    else:
        rendered = f"${amount:,.0f}"
    return f"{prefix}{rendered}".replace(".0K", "K").replace(".00M", "M")


def percent(value: float) -> str:
    return f"{round((value or 0) * 100)}%"


prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)


def blank_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_background(slide, color=WARM_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_box(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(1)):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    if fill_color is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    return shape


def add_text(slide, content, left, top, width, height, size=Pt(18), bold=False, color=DARK, align=PP_ALIGN.LEFT):
    text_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = text_box.text_frame
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = content
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return text_box


def add_bullets(slide, items, left, top, width, height, size=Pt(14), text_color=DARK, bullet_color=TEAL, gap=Pt(6)):
    text_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = text_box.text_frame
    frame.word_wrap = True
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.space_before = gap
        paragraph.clear()
        bullet = paragraph.add_run()
        bullet.text = "- "
        bullet.font.size = size
        bullet.font.bold = True
        bullet.font.color.rgb = bullet_color
        bullet.font.name = "Calibri"
        text_run = paragraph.add_run()
        text_run.text = item
        text_run.font.size = size
        text_run.font.color.rgb = text_color
        text_run.font.name = "Calibri"
    return text_box


def add_pill(slide, label, left, top, width=3.1):
    pill = add_box(slide, left, top, width, 0.34, fill_color=LIGHT_TEAL)
    frame = pill.text_frame
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = label.upper()
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = TEAL
    run.font.name = "Calibri"
    return pill


def add_top_bar(slide):
    add_box(slide, 0, 0, 13.33, 0.06, fill_color=TEAL)


def add_slide_number(slide, number):
    add_text(slide, f"{number}/{TOTAL_SLIDES}", 12.35, 7.05, 0.7, 0.3, size=Pt(9), color=MUTED, align=PP_ALIGN.RIGHT)


def add_metric_card(slide, left, top, width, value, label, tone="neutral", sublabel=None):
    value_color = {
        "alert": ALERT,
        "warning": WARNING,
        "accent": TEAL,
        "positive": GREEN,
        "neutral": DARK,
    }.get(tone, DARK)
    add_box(slide, left, top, width, 1.15, fill_color=LIGHT_CARD, line_color=LIGHT_BORDER)
    add_text(slide, value, left + 0.14, top + 0.10, width - 0.28, 0.42, size=Pt(20), bold=True, color=value_color)
    add_text(slide, label, left + 0.14, top + 0.58, width - 0.28, 0.22, size=Pt(9), bold=True, color=MUTED)
    if sublabel:
        add_text(slide, sublabel, left + 0.14, top + 0.82, width - 0.28, 0.18, size=Pt(8), color=MUTED)


def add_agent_card(slide, left, top, width, title, status, metric_line, tone="accent"):
    border_color = TEAL if tone == "accent" else LIGHT_BORDER
    fill_color = LIGHT_TEAL if tone == "accent" else LIGHT_CARD
    add_box(slide, left, top, width, 1.6, fill_color=fill_color, line_color=border_color, line_width=Pt(1.3))
    add_text(slide, title, left + 0.1, top + 0.12, width - 0.2, 0.58, size=Pt(12), bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_text(slide, status, left + 0.1, top + 0.78, width - 0.2, 0.22, size=Pt(9), bold=True, color=GREEN if status == "LIVE" else MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, metric_line, left + 0.1, top + 1.02, width - 0.2, 0.42, size=Pt(9), color=MUTED, align=PP_ALIGN.CENTER)


def build_deck(data: dict):
    summary = data["summary"]
    cases = data["cases"]
    predictions = data["predictions"]
    accounts = data["accounts"]
    drift = data["drift"]
    amendments = data["amendments"]
    billing = data["billing"]
    pricing = data["pricing"]
    ai_status = data["ai_status"]

    top_case = max(cases, key=lambda item: item["estimated_impact"])
    ordered_cases = sorted(cases, key=lambda item: item["estimated_impact"], reverse=True)
    top_prediction = predictions[0] if predictions else None
    top_drift = max(drift["analyses"], key=lambda item: item["total_estimated_annual_impact"])
    other_drifts = sorted(drift["analyses"], key=lambda item: item["total_estimated_annual_impact"], reverse=True)[1:]
    amendment_details = amendments["analyses"]
    top_amend_gain = max(amendment_details, key=lambda item: item["analysis"]["total_annual_revenue_delta"])
    top_amend_loss = min(amendment_details, key=lambda item: item["analysis"]["total_annual_revenue_delta"])
    flagged_billing = [item for item in billing["analyses"] if item["analysis"]["total_findings"] > 0 or item["analysis"]["status"] == "mismatch_detected"]
    top_billing = max(flagged_billing, key=lambda item: max(item["analysis"]["total_underbilled_amount"], item["analysis"]["total_overbilled_amount"]))
    pricing_ranked = sorted(pricing["recommendations"], key=lambda item: item["incremental_annual_contract_value"], reverse=True)
    top_pricing = pricing_ranked[0]

    queue_count = len(cases) + len(predictions) + len(flagged_billing)
    coverage_count = len(accounts)
    ai_mode = "AI enabled" if ai_status["enabled"] else "Rule-based fallback"

    # Slide 1
    slide = blank_slide()
    set_background(slide)
    add_top_bar(slide)
    add_box(slide, 0, 0.06, 5.05, 7.44, fill_color=TEAL)
    add_text(slide, "Commercial Execution\nIntelligence Platform", 0.45, 1.35, 4.3, 1.5, size=Pt(30), bold=True, color=WHITE)
    add_text(slide, "Current product deck\n5 live workspaces across pricing and post-sign execution", 0.45, 3.55, 4.15, 1.05, size=Pt(15), color=WHITE)
    add_text(slide, "June 2026", 0.45, 6.65, 2.0, 0.3, size=Pt(10), color=WHITE)

    add_text(slide, "AI-assisted commercial execution platform with a shared evidence layer and current live seeded metrics.", 5.55, 1.35, 7.0, 0.7, size=Pt(18), color=DARK)
    add_metric_card(slide, 5.55, 2.45, 3.15, fmt_currency(summary["total_estimated_missed_revenue"]), "Missed revenue found", tone="alert", sublabel="3 leakage cases")
    add_metric_card(slide, 8.95, 2.45, 3.15, compact_currency(drift["total_estimated_impact"]), "Quote-to-contract drift impact", tone="alert", sublabel="9 findings across 3 quotes")
    add_metric_card(slide, 5.55, 3.95, 3.15, compact_currency(amendments["net_annual_revenue_delta"], signed=True), "Net amendment delta", tone="positive", sublabel="11 impacts and 14 open actions")
    add_metric_card(slide, 8.95, 3.95, 3.15, compact_currency(pricing["total_incremental_annual_contract_value"], signed=True), "Pre-sign ACV upside", tone="positive", sublabel="7 uplift opportunities")
    add_box(slide, 5.55, 5.55, 6.55, 1.1, fill_color=LIGHT_CARD, line_color=LIGHT_BORDER)
    add_text(slide, f"Unified revenue integrity queue: {queue_count} active items across leakage, renewal risk, and billing audit.", 5.75, 5.75, 6.15, 0.3, size=Pt(12), bold=True, color=DARK)
    add_text(slide, f"Current demo mode: {ai_mode}. The same workflows can call AI services when enabled.", 5.75, 6.18, 6.15, 0.24, size=Pt(10), color=MUTED)
    add_slide_number(slide, 1)

    # Slide 2
    slide = blank_slide()
    set_background(slide)
    add_top_bar(slide)
    add_pill(slide, "The Problem", 0.6, 0.38)
    add_text(slide, "Revenue gets lost before and after signature.", 0.6, 0.82, 8.8, 0.7, size=Pt(28), bold=True, color=DARK)
    add_text(slide, "Pricing, contract truth, amendments, billing, and renewal actions drift unless one system keeps the evidence chain intact.", 0.6, 1.72, 9.6, 0.5, size=Pt(15), color=MUTED)
    add_bullets(slide, [
        "Final price is often set without same-company evidence, leaving ACV on the table.",
        "Quoted scope and renewal protections change before signature without a clean comparison loop.",
        "Amendments change commercial terms, but downstream systems keep old assumptions.",
        "Billing rate cards lag controlling contract economics even when the contract is clear.",
        "Renewal notice windows pass and recurring revenue rights expire silently.",
    ], 0.6, 2.55, 7.6, 3.5)
    add_box(slide, 8.8, 1.75, 4.0, 4.95, fill_color=LIGHT_TEAL, line_color=TEAL, line_width=Pt(1.4))
    add_text(slide, "Current proof in the demo data", 9.0, 1.95, 3.6, 0.25, size=Pt(10), bold=True, color=TEAL)
    add_metric_card(slide, 9.0, 2.35, 1.7, str(queue_count), "Active revenue integrity items", tone="neutral")
    add_metric_card(slide, 10.95, 2.35, 1.7, str(drift["total_findings"]), "Drift findings", tone="alert")
    add_metric_card(slide, 9.0, 3.75, 1.7, str(amendments["total_impacts"]), "Amendment impacts", tone="warning")
    add_metric_card(slide, 10.95, 3.75, 1.7, str(pricing["opportunities_with_price_uplift"]), "Pricing lift opportunities", tone="positive")
    add_text(slide, "The same evidence layer has to support pre-sign pricing decisions, post-sign drift review, amendment routing, billing audit, and renewal protection.", 9.0, 5.2, 3.55, 0.9, size=Pt(12), color=DARK)
    add_slide_number(slide, 2)

    # Slide 3
    slide = blank_slide()
    set_background(slide)
    add_top_bar(slide)
    add_pill(slide, "Platform Architecture", 0.6, 0.38)
    add_text(slide, "One evidence layer, five live workspaces.", 0.6, 0.82, 9.0, 0.7, size=Pt(28), bold=True, color=DARK)
    add_box(slide, 1.0, 2.2, 11.3, 0.78, fill_color=TEAL)
    add_text(slide, "Shared evidence layer | contracts | amendments | invoices | quotes | renewal events | same-company deal history", 1.1, 2.38, 11.1, 0.3, size=Pt(11), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_agent_card(slide, 0.9, 3.35, 2.35, "Revenue Integrity\nWorkspace", "LIVE", f"{fmt_currency(summary['total_estimated_missed_revenue'])} missed | {fmt_currency(summary['total_predicted_at_risk_revenue'])} risk", tone="accent")
    add_agent_card(slide, 3.45, 3.35, 2.35, "Quote-to-Contract\nDrift Detector", "LIVE", f"{compact_currency(drift['total_estimated_impact'])} | {drift['total_findings']} findings", tone="accent")
    add_agent_card(slide, 6.0, 3.35, 2.35, "Amendment Impact\nDetector", "LIVE", f"{compact_currency(amendments['net_annual_revenue_delta'], signed=True)} | {amendments['total_action_items_open']} actions", tone="accent")
    add_agent_card(slide, 8.55, 3.35, 2.35, "Billing vs Contract\nMismatch", "LIVE", f"{fmt_currency(billing['total_underbilled_amount'])} underbilled", tone="accent")
    add_agent_card(slide, 11.1, 3.35, 1.3, "Pricing\nAdvisor", "LIVE", compact_currency(pricing['total_incremental_annual_contract_value'], signed=True), tone="accent")
    add_box(slide, 1.0, 5.45, 11.3, 0.95, fill_color=LIGHT_CARD, line_color=LIGHT_BORDER)
    add_text(slide, f"Current operating mode: {ai_mode}. Extraction and explanation fall back deterministically in this environment while the same product surface remains AI-ready.", 1.2, 5.72, 10.9, 0.3, size=Pt(12), color=DARK, align=PP_ALIGN.CENTER)
    add_slide_number(slide, 3)

    # Slide 4
    slide = blank_slide()
    set_background(slide)
    add_top_bar(slide)
    add_pill(slide, "Unified Revenue Integrity", 0.6, 0.38, width=3.6)
    add_text(slide, "One operator view for leakage, renewal risk, billing, and account coverage.", 0.6, 0.82, 10.6, 0.8, size=Pt(26), bold=True, color=DARK)
    add_bullets(slide, [
        "Unified queue combines leakage cases, renewal risk, and flagged billing audits.",
        "All Accounts coverage keeps low-signal accounts visible instead of disappearing from the workflow.",
        "Classic leakage and billing pages still exist for deep investigation when needed.",
        "Document uploads, evidence review, and source-of-record detail stay in the same workspace.",
    ], 0.6, 2.0, 7.2, 3.0)
    add_metric_card(slide, 8.4, 1.95, 1.9, str(len(cases)), "Leakage cases", tone="alert", sublabel=fmt_currency(summary["total_estimated_missed_revenue"]))
    add_metric_card(slide, 10.45, 1.95, 1.9, str(len(predictions)), "Renewal risks", tone="warning", sublabel=fmt_currency(summary["total_predicted_at_risk_revenue"]))
    add_metric_card(slide, 8.4, 3.35, 1.9, str(len(flagged_billing)), "Billing audits", tone="accent", sublabel=fmt_currency(billing["total_underbilled_amount"]))
    add_metric_card(slide, 10.45, 3.35, 1.9, str(coverage_count), "Accounts covered", tone="neutral", sublabel="All accounts queue")
    add_box(slide, 0.6, 5.35, 12.1, 1.05, fill_color=LIGHT_CARD, line_color=LIGHT_BORDER)
    add_text(slide, f"Live examples: {top_case['account_name']} shows {fmt_currency(top_case['estimated_impact'])} in missed uplift. {top_prediction['account_name']} has {top_prediction['days_until_deadline']} days left before {fmt_currency(top_prediction['predicted_impact'])} is at risk.", 0.85, 5.63, 11.6, 0.35, size=Pt(12), color=DARK)
    add_slide_number(slide, 4)

    # Slide 5
    slide = blank_slide()
    set_background(slide)
    add_top_bar(slide)
    add_pill(slide, "Agent 1", 0.6, 0.38)
    add_text(slide, "Revenue Leakage Investigator", 0.6, 0.82, 7.6, 0.6, size=Pt(28), bold=True, color=TEAL)
    add_text(slide, "Missed revenue and preventive renewal risk are already quantified from the contract dossier.", 0.6, 1.62, 8.9, 0.45, size=Pt(15), color=MUTED)
    add_bullets(slide, [
        "Resolves the controlling commercial term across PDFs and DOCX files.",
        "Quantifies missed uplift and near-term renewal risk with deterministic math.",
        "Explains which document controls and what action operations should take next.",
        "Supports retrospective leakage detection and preventive intervention.",
    ], 0.6, 2.35, 6.7, 2.9)
    add_box(slide, 8.15, 2.1, 4.4, 4.5, fill_color=LIGHT_TEAL, line_color=TEAL, line_width=Pt(1.4))
    add_text(slide, "Current live proof", 8.35, 2.28, 4.0, 0.22, size=Pt(10), bold=True, color=TEAL)
    add_text(slide, fmt_currency(summary["total_estimated_missed_revenue"]), 8.35, 2.7, 1.8, 0.4, size=Pt(22), bold=True, color=ALERT)
    add_text(slide, "missed revenue across 3 cases", 10.0, 2.78, 2.0, 0.25, size=Pt(10), color=MUTED)
    add_text(slide, fmt_currency(summary["total_predicted_at_risk_revenue"]), 8.35, 3.35, 1.8, 0.4, size=Pt(22), bold=True, color=WARNING)
    add_text(slide, "at risk across 1 preventive case", 10.0, 3.43, 2.2, 0.25, size=Pt(10), color=MUTED)
    add_text(slide, f"Top case: {top_case['account_name']}\n{fmt_currency(top_case['estimated_impact'])} missed because an {round(((top_case['expected_value'] / top_case['actual_value']) - 1) * 100)}% uplift never hit billing.", 8.35, 4.05, 3.9, 0.9, size=Pt(12), color=DARK)
    add_text(slide, f"Other open cases: {ordered_cases[1]['account_name']} {fmt_currency(ordered_cases[1]['estimated_impact'])} and {ordered_cases[2]['account_name']} {fmt_currency(ordered_cases[2]['estimated_impact'])}.", 8.35, 5.15, 3.9, 0.5, size=Pt(11), color=DARK)
    add_text(slide, f"Preventive example: {top_prediction['account_name']} must send the uplift notice within {top_prediction['days_until_deadline']} days.", 8.35, 5.8, 3.9, 0.45, size=Pt(11), color=ALERT)
    add_slide_number(slide, 5)

    # Slide 6
    slide = blank_slide()
    set_background(slide)
    add_top_bar(slide)
    add_pill(slide, "Agent 2", 0.6, 0.38)
    add_text(slide, "Quote-to-Contract Drift Detector", 0.6, 0.82, 8.5, 0.6, size=Pt(28), bold=True, color=TEAL)
    add_text(slide, "Signed terms are compared to approved quotes before pricing concessions quietly become permanent.", 0.6, 1.62, 9.3, 0.45, size=Pt(15), color=MUTED)
    add_metric_card(slide, 0.6, 2.3, 2.45, str(drift["total_quotes_analyzed"]), "Quotes analyzed", tone="neutral")
    add_metric_card(slide, 3.25, 2.3, 2.45, str(drift["total_findings"]), "Drift findings", tone="warning")
    add_metric_card(slide, 5.9, 2.3, 2.45, str(drift["total_high_severity"]), "High severity", tone="alert")
    add_metric_card(slide, 8.55, 2.3, 3.0, compact_currency(drift["total_estimated_impact"]), "Estimated annual impact", tone="alert")
    add_bullets(slide, [
        f"{top_drift['account_name']}: {compact_currency(top_drift['total_estimated_annual_impact'])} annual exposure from scope and renewal drift.",
        f"{other_drifts[0]['account_name']}: {compact_currency(other_drifts[0]['total_estimated_annual_impact'])} impact from renewal caps and support downgrades.",
        f"{other_drifts[1]['account_name']}: {compact_currency(other_drifts[1]['total_estimated_annual_impact'])} impact from seat-count reduction.",
        "Findings cover scope, quantity, support tier, and renewal term changes.",
    ], 0.6, 4.0, 12.0, 2.1)
    add_slide_number(slide, 6)

    # Slide 7
    slide = blank_slide()
    set_background(slide)
    add_top_bar(slide)
    add_pill(slide, "Agent 3", 0.6, 0.38)
    add_text(slide, "Amendment Impact Detector", 0.6, 0.82, 7.6, 0.6, size=Pt(28), bold=True, color=TEAL)
    add_text(slide, "Amendments trigger operational action only when before/after changes are translated into billing, provisioning, and workflow updates.", 0.6, 1.62, 10.7, 0.55, size=Pt(15), color=MUTED)
    add_metric_card(slide, 0.6, 2.25, 2.2, str(amendments["total_analyses"]), "Amendments analyzed", tone="neutral")
    add_metric_card(slide, 2.95, 2.25, 2.2, str(amendments["total_impacts"]), "Detected impacts", tone="warning")
    add_metric_card(slide, 5.3, 2.25, 2.2, str(amendments["total_action_items_open"]), "Open actions", tone="alert")
    add_metric_card(slide, 7.65, 2.25, 2.75, compact_currency(amendments["net_annual_revenue_delta"], signed=True), "Net annual delta", tone="positive")
    add_metric_card(slide, 10.6, 2.25, 1.7, f"{amendments['positive_amendments']}/{amendments['negative_amendments']}", "Positive / negative", tone="neutral")
    add_box(slide, 0.6, 4.0, 5.95, 2.0, fill_color=LIGHT_CARD, line_color=LIGHT_BORDER)
    add_text(slide, f"Expansion upside\n{top_amend_gain['analysis']['account_name']} adds {compact_currency(top_amend_gain['analysis']['total_annual_revenue_delta'], signed=True)} in annual revenue.", 0.85, 4.25, 5.4, 0.6, size=Pt(16), bold=True, color=GREEN)
    add_text(slide, top_amend_gain["analysis"]["amendment_summary"], 0.85, 4.95, 5.4, 0.65, size=Pt(11), color=DARK)
    add_box(slide, 6.8, 4.0, 5.95, 2.0, fill_color=LIGHT_CARD, line_color=LIGHT_BORDER)
    add_text(slide, f"Contraction risk\n{top_amend_loss['analysis']['account_name']} removes {compact_currency(top_amend_loss['analysis']['total_annual_revenue_delta'], signed=True)} in annual revenue.", 7.05, 4.25, 5.4, 0.6, size=Pt(16), bold=True, color=ALERT)
    add_text(slide, top_amend_loss["analysis"]["amendment_summary"], 7.05, 4.95, 5.4, 0.65, size=Pt(11), color=DARK)
    add_slide_number(slide, 7)

    # Slide 8
    slide = blank_slide()
    set_background(slide)
    add_top_bar(slide)
    add_pill(slide, "Agents 4 and 5", 0.6, 0.38)
    add_text(slide, "Billing mismatch and pre-sign pricing already sit on the same platform.", 0.6, 0.82, 11.8, 0.7, size=Pt(26), bold=True, color=DARK)
    add_box(slide, 0.6, 1.9, 6.0, 4.7, fill_color=LIGHT_CARD, line_color=LIGHT_BORDER)
    add_text(slide, "Billing vs Contract Mismatch", 0.9, 2.15, 5.4, 0.35, size=Pt(18), bold=True, color=TEAL)
    add_bullets(slide, [
        f"{billing['total_contracts_monitored']} contracts monitored, {billing['flagged_contracts']} flagged.",
        f"{billing['total_findings']} findings with {billing['high_severity_findings']} high severity.",
        f"{fmt_currency(billing['total_underbilled_amount'])} underbilled and {fmt_currency(billing['total_overbilled_amount'])} overbilled.",
        f"Top mismatch: {top_billing['analysis']['account_name']} missed {fmt_currency(top_billing['analysis']['total_underbilled_amount'])} across {top_billing['analysis']['total_findings']} invoices after the uplift should have applied.",
    ], 0.9, 2.75, 5.2, 2.3)
    add_text(slide, "This confirms the same contract truth can drive invoice-level audit, not just account-level leakage scoring.", 0.9, 5.45, 5.1, 0.5, size=Pt(11), color=MUTED)

    add_box(slide, 6.75, 1.9, 6.0, 4.7, fill_color=LIGHT_CARD, line_color=LIGHT_BORDER)
    add_text(slide, "Pre-Sign Pricing Advisor", 7.05, 2.15, 5.4, 0.35, size=Pt(18), bold=True, color=TEAL)
    add_bullets(slide, [
        f"{pricing['total_opportunities']} opportunities scored and {pricing['opportunities_with_price_uplift']} uplift recommendations.",
        f"{compact_currency(pricing['total_incremental_annual_contract_value'], signed=True)} in incremental annual contract value.",
        f"Average improved close confidence: {percent(pricing['average_improved_close_confidence'])}.",
        f"Top recommendation: {top_pricing['account_name']} at {compact_currency(top_pricing['incremental_annual_contract_value'], signed=True)} annual lift with {percent(top_pricing['improved_price_close_confidence'])} projected close confidence.",
    ], 7.05, 2.75, 5.2, 2.3)
    add_text(slide, f"Next best examples: {pricing_ranked[1]['account_name']} {compact_currency(pricing_ranked[1]['incremental_annual_contract_value'], signed=True)} and {pricing_ranked[2]['account_name']} {compact_currency(pricing_ranked[2]['incremental_annual_contract_value'], signed=True)}.", 7.05, 5.45, 5.2, 0.5, size=Pt(11), color=MUTED)
    add_slide_number(slide, 8)

    # Slide 9
    slide = blank_slide()
    set_background(slide)
    add_top_bar(slide)
    add_pill(slide, "Close", 0.6, 0.38)
    add_text(slide, "Commercial execution intelligence is now a product surface, not a one-off demo.", 0.6, 0.82, 11.8, 0.8, size=Pt(27), bold=True, color=DARK)
    add_bullets(slide, [
        "Five live workspaces already reuse one evidence layer instead of rebuilding ingestion and resolution each time.",
        f"The current environment proves value even in {ai_mode.lower()}, which keeps rollout conservative and deterministic.",
        "The same core can extend into obligation tracking, renewal playbooks, controls, and operator copilots without changing the data model.",
        "Contracts become the commercial source of truth only when pricing, billing, renewal, and operations all act on the same evidence.",
    ], 0.8, 2.0, 8.8, 2.9, size=Pt(16))
    add_box(slide, 9.8, 1.95, 2.7, 3.9, fill_color=LIGHT_TEAL, line_color=TEAL, line_width=Pt(1.4))
    add_text(slide, "Live now", 10.15, 2.15, 2.0, 0.22, size=Pt(10), bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    add_text(slide, "Revenue integrity\nDrift detection\nAmendment impact\nBilling mismatch\nPricing advisor", 10.15, 2.65, 2.0, 1.6, size=Pt(14), bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_text(slide, "Closing line", 10.15, 4.65, 2.0, 0.22, size=Pt(10), bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    add_text(slide, "Turn contract truth into pricing, revenue integrity, and operational action.", 10.05, 5.0, 2.2, 0.65, size=Pt(12), color=DARK, align=PP_ALIGN.CENTER)
    add_slide_number(slide, 9)


def main():
    data = {
        "summary": fetch_json("/api/dashboard/summary"),
        "cases": ensure_list(fetch_json("/api/cases")),
        "predictions": ensure_list(fetch_json("/api/predictions")),
        "accounts": ensure_list(fetch_json("/api/accounts")),
        "drift": fetch_json("/api/drift/dashboard"),
        "amendments": fetch_json("/api/amendments/dashboard"),
        "billing": fetch_json("/api/billing-mismatch/dashboard"),
        "pricing": fetch_json("/api/pricing/dashboard"),
        "ai_status": fetch_json("/api/system/ai-status"),
    }

    build_deck(data)
    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(f"Created {OUTPUT_PATH}")


if __name__ == "__main__":
    main()